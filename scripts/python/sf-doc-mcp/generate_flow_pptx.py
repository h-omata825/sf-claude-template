#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
業務フロー図 PPTX 生成スクリプト（Mermaid sequenceDiagram 版）

swimlanes.json の各フローを Mermaid sequenceDiagram → PNG → PPTX スライドに変換する。

入力:
  docs/flow/swimlanes.json

出力:
  業務フロー図.pptx

Usage:
  python generate_flow_pptx.py \\
    --docs-dir <path/to/project/docs> \\
    --output-dir <output dir> \\
    --author "作成者名"
"""
import argparse
import datetime
import json
import shutil
import subprocess
import sys
import tempfile
from collections import Counter
from itertools import permutations
from pathlib import Path

from tmp_utils import get_project_tmp_dir, set_project_tmp_dir

# ── 定数 ────────────────────────────────────────────────────────────────────

SLIDE_W_IN = 13.33   # 16:9 スライド幅 (インチ)
SLIDE_H_IN = 7.5     # スライド高さ
TITLE_H_IN = 0.55    # タイトル行の高さ
IMG_PAD    = 0.1     # 画像の左右パディング

MERMAID_INIT = (
    "%%{init: {'theme': 'base', 'themeVariables': {"
    "'fontSize': '13px', "
    "'actorBkg': '#5B6FA6', 'actorBorder': '#3d4f80', "
    "'actorTextColor': '#ffffff', "
    "'signalTextColor': '#222222', "
    "'noteBkgColor': '#fffff0', 'noteTextColor': '#444444', "
    "'sequenceNumberColor': '#ffffff'"
    "}}}%%"
)

PHASE_COLORS = [
    "rgb(220,235,255)",  # 青
    "rgb(220,255,220)",  # 緑
    "rgb(255,245,215)",  # 黄
    "rgb(255,228,215)",  # オレンジ
    "rgb(238,218,255)",  # 紫
    "rgb(215,245,245)",  # 水色
]

# ── 参加者名の短縮 ────────────────────────────────────────────────────────────

def _shorten_participant(name: str, max_len: int = 10) -> str:
    """括弧内を除去して短縮する。"""
    import re
    short = re.sub(r'[\(（][^)）]*[\)）]', '', name).strip()
    return short if short else name


def _build_aliases(participants: list[str]) -> dict[str, str]:
    """参加者名 → 表示名（重複しない短縮名）を返す。"""
    aliases: dict[str, str] = {}
    used: set[str] = set()
    for p in participants:
        short = _shorten_participant(p)
        if short in used:
            short = p  # 重複時はフル名を使う
        aliases[p] = short
        used.add(short)
    return aliases


# ── 参加者の最適配置 ─────────────────────────────────────────────────────────

def _optimal_order(participants: list[str], pair_freq: Counter) -> list[str]:
    """矢印の交差長の合計を最小化する参加者順序を総当たりで求める (n≦8 向け)。"""
    if len(participants) <= 1:
        return participants
    best_order, best_cost = None, float('inf')
    for perm in permutations(participants):
        pos = {p: i for i, p in enumerate(perm)}
        cost = sum(
            freq * abs(pos.get(a, 0) - pos.get(b, 0))
            for (a, b), freq in pair_freq.items()
        )
        if cost < best_cost:
            best_cost = cost
            best_order = list(perm)
    return best_order


def _calc_pair_freq(flow: dict) -> Counter:
    steps = {s['id']: s['lane'] for s in flow['steps']}
    freq: Counter = Counter()
    for t in flow.get('transitions', []):
        a = steps.get(t['from'], '')
        b = steps.get(t['to'], '')
        if a and b and a != b:
            freq[tuple(sorted([a, b]))] += 1
    return freq


# ── フロー → Mermaid コード ──────────────────────────────────────────────────

def _flow_to_mermaid(flow: dict, phases: list[dict] | None = None) -> str:
    """flow dict を Mermaid sequenceDiagram テキストに変換する。"""
    steps  = {s['id']: s for s in flow['steps']}
    sorted_ids = sorted(steps.keys())

    # 遷移マップ
    to_map: dict[int, list[tuple[int, str]]] = {}
    for t in flow.get('transitions', []):
        to_map.setdefault(t['to'], []).append((t['from'], t.get('condition', '')))

    # 参加者の最適順序
    participants = [l['name'] for l in flow['lanes']]
    pair_freq = _calc_pair_freq(flow)
    ordered = _optimal_order(participants, pair_freq)
    aliases = _build_aliases(ordered)

    lines = [MERMAID_INIT, 'sequenceDiagram']
    for p in ordered:
        a = aliases[p]
        lines.append(f'  participant {a}')
    lines.append('')

    # フェーズ分割がある場合は rect ブロックで囲む
    if phases:
        phase_step_sets = [set(ph['step_ids']) for ph in phases]
    else:
        phase_step_sets = None

    open_phase = -1
    prev_lane: str | None = None

    for sid in sorted_ids:
        s = steps[sid]
        lane  = aliases[s['lane']]
        label = s['title']

        # このステップへの遷移に条件があるか（任意フロー）
        incoming = to_map.get(sid, [])
        is_optional = any(cond for _, cond in incoming)

        # フェーズ開始/終了
        if phase_step_sets:
            for pi, pset in enumerate(phase_step_sets):
                if sid in pset and pi != open_phase:
                    if open_phase >= 0:
                        lines.append('  end')
                    color = PHASE_COLORS[pi % len(PHASE_COLORS)]
                    ph_label = phases[pi]['title']
                    all_aliases = list(dict.fromkeys(aliases.values()))
                    left, right = all_aliases[0], all_aliases[-1]
                    lines.append(f'  rect {color}')
                    lines.append(f'    Note over {left},{right}: {ph_label}')
                    open_phase = pi
                    break

        # 矢印の種別: 任意フローは破線
        arrow = '-->>' if is_optional else '->>'

        if prev_lane is None or prev_lane == lane:
            lines.append(f'  Note over {lane}: {label}')
        else:
            lines.append(f'  {prev_lane}{arrow}{lane}: {label}')

        prev_lane = lane

    if open_phase >= 0:
        lines.append('  end')

    return '\n'.join(lines)


# ── Mermaid → PNG ────────────────────────────────────────────────────────────

def _render_mermaid(mmd_text: str, output_png: Path, width: int = 1600) -> bool:
    """mmdc (mermaid CLI) で PNG を生成する。失敗時は False を返す。"""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.mmd',
                                    encoding='utf-8', delete=False,
                                    dir=get_project_tmp_dir()) as f:
        f.write(mmd_text)
        mmd_path = Path(f.name)
    # Windows では npx.cmd を使う必要がある
    npx_cmd = 'npx.cmd' if sys.platform == 'win32' else 'npx'
    try:
        result = subprocess.run(
            [npx_cmd, '@mermaid-js/mermaid-cli',
             '-i', str(mmd_path),
             '-o', str(output_png),
             '-w', str(width),
             '--backgroundColor', 'white'],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            print(f'[WARN] mmdc error: {result.stderr[:200]}', file=sys.stderr)
            return False
        return True
    except Exception as e:
        print(f'[WARN] mmdc failed: {e}', file=sys.stderr)
        return False
    finally:
        mmd_path.unlink(missing_ok=True)


# ── PNG → PPTX ───────────────────────────────────────────────────────────────

def _build_pptx(
    slides_data: list[dict],   # [{'title': str, 'png': Path}]
    output_path: Path,
    author: str,
    company: str,
    date_str: str,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from PIL import Image

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]  # blank

    # ── 表紙スライド ──
    cover = prs.slides.add_slide(blank_layout)
    # 背景色
    bg = cover.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x3D, 0x4F, 0x80)

    def _add_text(slide, text, left, top, width, height, size, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF), align=None):
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if align:
            p.alignment = align
        return txb

    from pptx.enum.text import PP_ALIGN
    _add_text(cover, 'シーケンス図', 1, 2.5, 11.33, 1.2, 32, bold=True, align=PP_ALIGN.CENTER)
    _add_text(cover, company or '', 1, 3.9, 11.33, 0.6, 16, align=PP_ALIGN.CENTER)
    _add_text(cover, f'{date_str}　作成: {author}', 1, 5.5, 11.33, 0.5, 12,
              color=RGBColor(0xCC, 0xD6, 0xFF), align=PP_ALIGN.CENTER)

    # ── フロースライド ──
    for item in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        png: Path = item['png']
        title: str = item['title']

        # タイトルバー
        title_box = slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(TITLE_H_IN))
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        fill2 = title_box.fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF8)

        # 画像を縦横比を保ってスライドに収める
        avail_w = SLIDE_W_IN - IMG_PAD * 2
        avail_h = SLIDE_H_IN - TITLE_H_IN - 0.05
        try:
            with Image.open(png) as im:
                img_w, img_h = im.size
            ratio = img_w / img_h
            fit_w = avail_w
            fit_h = fit_w / ratio
            if fit_h > avail_h:
                fit_h = avail_h
                fit_w = fit_h * ratio
            left = (SLIDE_W_IN - fit_w) / 2
            top  = TITLE_H_IN + 0.05
            slide.shapes.add_picture(str(png), Inches(left), Inches(top),
                                     Inches(fit_w), Inches(fit_h))
        except Exception as e:
            print(f'[WARN] 画像挿入失敗: {png} — {e}', file=sys.stderr)

    prs.save(str(output_path))


# ── フェーズ自動検出 ─────────────────────────────────────────────────────────

def _auto_phases(flow: dict, steps_per_phase: int = 5) -> list[dict] | None:
    """ステップ数が多い場合のみ自動フェーズ分割を行う。"""
    steps = sorted(flow['steps'], key=lambda s: s['id'])
    if len(steps) <= 8:
        return None  # 短いフローはフェーズ不要

    phases = []
    for i in range(0, len(steps), steps_per_phase):
        chunk = steps[i:i + steps_per_phase]
        first, last = chunk[0]['id'], chunk[-1]['id']
        phases.append({
            'title': f'フェーズ {i // steps_per_phase + 1}（Step {first}〜{last}）',
            'step_ids': {s['id'] for s in chunk},
        })
    return phases


# ── エントリポイント ─────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(description='業務フロー図 PPTX 生成（Mermaid版）')
    ap.add_argument('--docs-dir',   required=True)
    ap.add_argument('--output-dir', required=True)
    ap.add_argument('--author',     default='')
    ap.add_argument('--width',      type=int, default=1600, help='PNG 生成幅 (px)')
    args = ap.parse_args()

    docs_dir   = Path(args.docs_dir)
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    set_project_tmp_dir(output_dir)

    swimlane_path = docs_dir / 'flow' / 'swimlanes.json'
    if not swimlane_path.exists():
        print(f'ERROR: {swimlane_path} が見つかりません', file=sys.stderr)
        sys.exit(1)

    data = json.loads(swimlane_path.read_text(encoding='utf-8'))
    flows = data.get('flows', [])
    if not flows:
        print('ERROR: flows が空です', file=sys.stderr)
        sys.exit(1)

    # 会社名を org-profile.md から取得
    import re
    company = ''
    profile = docs_dir / 'overview' / 'org-profile.md'
    if profile.exists():
        m = re.search(r'\|\s*会社名\s*\|\s*(.+?)\s*\|', profile.read_text(encoding='utf-8'))
        if m:
            company = m.group(1).strip()

    date_str = datetime.date.today().strftime('%Y年%m月')

    slides_data: list[dict] = []

    with tempfile.TemporaryDirectory(dir=get_project_tmp_dir()) as tmp:
        tmp_dir = Path(tmp)
        for i, flow in enumerate(flows):
            flow_id    = flow.get('id', f'flow{i}')
            flow_title = flow.get('title', f'フロー {i+1}')
            print(f'  [{i+1}/{len(flows)}] {flow_title} を生成中...', file=sys.stderr)

            phases = _auto_phases(flow)
            mmd = _flow_to_mermaid(flow, phases)

            png_path = tmp_dir / f'{flow_id}.png'
            ok = _render_mermaid(mmd, png_path, width=args.width)
            if not ok or not png_path.exists():
                print(f'  [SKIP] {flow_title}: PNG 生成失敗', file=sys.stderr)
                continue

            slides_data.append({'title': flow_title, 'png': png_path})

        if not slides_data:
            print('ERROR: 生成できたスライドが 0 件です', file=sys.stderr)
            sys.exit(1)

        output_path = output_dir / 'シーケンス図.pptx'
        _build_pptx(slides_data, output_path, args.author, company, date_str)

    print(f'出力: {output_path}', file=sys.stderr)
    print(str(output_path))


if __name__ == '__main__':
    main()
