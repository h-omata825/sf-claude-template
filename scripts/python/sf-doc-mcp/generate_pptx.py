# -*- coding: utf-8 -*-
"""JSON → PowerPoint 生成スクリプト

Claude が再構成したスライド構成 JSON を受け取り、統一フォーマットの PowerPoint を生成する。
全資料タイプ（システム概要書・業務フロー等）で共通使用。

対応レイアウト:
  section     — セクション区切り（濃紺背景）
  content     — 本文テキスト
  table       — テーブル
  bullets     — 箇条書き
  two_column  — 2カラム
  toc         — 目次（セクションへのハイパーリンク付き）
  diagram     — 図（ボックス + 矢印 + グループ）
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


# ── カラーパレット ──
C = {
    "primary":    RGBColor(0x1E, 0x3A, 0x5F),   # 濃紺
    "secondary":  RGBColor(0x4A, 0x90, 0xD9),   # 青
    "accent":     RGBColor(0xE8, 0x6C, 0x00),   # オレンジ
    "bg_dark":    RGBColor(0x1E, 0x3A, 0x5F),   # 背景（濃紺）
    "bg_light":   RGBColor(0xF5, 0xF7, 0xFA),   # 背景（薄灰）
    "text_dark":  RGBColor(0x2D, 0x2D, 0x2D),   # 本文
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),   # 白文字
    "text_sub":   RGBColor(0x6B, 0x6B, 0x6B),   # サブテキスト
    "border":     RGBColor(0xD0, 0xD0, 0xD0),   # 罫線
    "table_hd":   RGBColor(0x2C, 0x3E, 0x50),   # テーブルヘッダー
    "table_alt":  RGBColor(0xEB, 0xF0, 0xF5),   # テーブル交互色
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
}

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(1.4)
CONTENT_W = Inches(11.7)
CONTENT_H = Inches(5.2)
FONT_NAME = "BIZ UDPゴシック"
FONT_NAME_EN = "Segoe UI"


def _swimlane_trunc(text: str, max_w: int = 22) -> str:
    """スイムレーンのステップラベルを表示幅でトリムする（全角=2、半角=1）"""
    import unicodedata
    total, out = 0, []
    for c in text:
        w = 2 if unicodedata.east_asian_width(c) in ('W', 'F') else 1
        if total + w > max_w:
            out.append('…')
            break
        out.append(c)
        total += w
    return ''.join(out)


def _add_bg_rect(slide, x, y, w, h, color, alpha=None):
    """背景矩形を追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def _add_text_box(slide, left, top, width, height):
    """テキストボックスを追加"""
    return slide.shapes.add_textbox(left, top, width, height)


def _set_paragraph(para, text, font_size=12, bold=False, color=None, align=None, font_name=None):
    """段落のフォーマットを設定"""
    para.text = text
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.name = font_name or FONT_NAME
    if color:
        para.font.color.rgb = color
    if align:
        para.alignment = align


def _add_footer(slide, page_num: int, total: int, title: str):
    """フッターを追加"""
    # ページ番号
    tb = _add_text_box(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4), Inches(1.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    _set_paragraph(p, f"{page_num} / {total}", font_size=8, color=C["text_sub"], align=PP_ALIGN.RIGHT)

    # Confidential
    tb2 = _add_text_box(slide, Inches(0.3), SLIDE_H - Inches(0.4), Inches(3), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    _set_paragraph(p2, "Confidential", font_size=8, color=C["text_sub"])


def build_cover(prs, data: dict):
    """表紙スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 背景
    _add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C["bg_dark"])

    # アクセントライン
    _add_bg_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.05), C["accent"])

    # タイトル
    tb = _add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    _set_paragraph(p, data.get("title", ""), font_size=36, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    # サブタイトル
    if data.get("subtitle"):
        tb2 = _add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]
        _set_paragraph(p2, data["subtitle"], font_size=18, color=C["secondary"], align=PP_ALIGN.CENTER)

    # メタ情報
    meta_lines = []
    if data.get("company"):
        meta_lines.append(data["company"])
    meta_lines.append(f"バージョン {data.get('version', '1.0')}  |  {data.get('date', '')}")
    if data.get("author"):
        meta_lines.append(f"作成者: {data['author']}")

    tb3 = _add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(1.5))
    tf = tb3.text_frame
    for i, line in enumerate(meta_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_paragraph(p, line, font_size=12, color=C["text_sub"], align=PP_ALIGN.CENTER)
        p.space_after = Pt(4)


def build_section(prs, slide_data: dict, page_num: int, total: int):
    """セクション区切りスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    _add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C["primary"])

    # セクション番号（あれば）
    tb = _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    _set_paragraph(p, slide_data.get("title", ""), font_size=32, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    if slide_data.get("body"):
        tb2 = _add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(1.0))
        p2 = tb2.text_frame.paragraphs[0]
        _set_paragraph(p2, slide_data["body"], font_size=14, color=C["text_light"], align=PP_ALIGN.CENTER)


def _add_slide_title(slide, title: str):
    """スライドタイトル帯"""
    # タイトル背景帯
    _add_bg_rect(slide, 0, 0, SLIDE_W, Inches(1.0), C["primary"])
    tb = _add_text_box(slide, MARGIN_LEFT, Inches(0.15), CONTENT_W, Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    _set_paragraph(p, title, font_size=20, bold=True, color=C["white"])


def build_content(prs, slide_data: dict, page_num: int, total: int):
    """本文スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    body = slide_data.get("body", "")
    if body:
        tb = _add_text_box(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, CONTENT_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_paragraph(p, line, font_size=12, color=C["text_dark"])
            p.space_after = Pt(6)

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


def build_bullets(prs, slide_data: dict, page_num: int, total: int):
    """箇条書きスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    bullets = slide_data.get("bullets", [])
    if not bullets:
        _add_footer(slide, page_num, total, slide_data.get("title", ""))
        return

    # 本文（あれば）
    y_start = MARGIN_TOP
    if slide_data.get("body"):
        tb_body = _add_text_box(slide, MARGIN_LEFT, y_start, CONTENT_W, Inches(0.8))
        p = tb_body.text_frame.paragraphs[0]
        _set_paragraph(p, slide_data["body"], font_size=12, color=C["text_dark"])
        y_start += Inches(0.9)

    tb = _add_text_box(slide, MARGIN_LEFT, y_start, CONTENT_W, SLIDE_H - y_start - Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, dict):
            # {title: "xxx", detail: "yyy"} 形式
            _set_paragraph(p, f"■ {item.get('title', '')}", font_size=12, bold=True, color=C["primary"])
            if item.get("detail"):
                p2 = tf.add_paragraph()
                _set_paragraph(p2, f"    {item['detail']}", font_size=11, color=C["text_dark"])
                p2.space_after = Pt(8)
        else:
            _set_paragraph(p, f"•  {item}", font_size=12, color=C["text_dark"])
            p.space_after = Pt(6)

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


def build_table(prs, slide_data: dict, page_num: int, total: int):
    """テーブルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    tbl_data = slide_data.get("table", {})
    headers = tbl_data.get("headers", [])
    rows = tbl_data.get("rows", [])

    if not headers:
        _add_footer(slide, page_num, total, slide_data.get("title", ""))
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    # テーブルサイズ
    tbl_w = CONTENT_W
    row_h = Inches(0.35)
    tbl_h = row_h * min(n_rows, 15)
    tbl_top = MARGIN_TOP

    if slide_data.get("body"):
        tb_body = _add_text_box(slide, MARGIN_LEFT, tbl_top, CONTENT_W, Inches(0.6))
        p = tb_body.text_frame.paragraphs[0]
        _set_paragraph(p, slide_data["body"], font_size=11, color=C["text_dark"])
        tbl_top += Inches(0.7)

    table_shape = slide.shapes.add_table(min(n_rows, 15), n_cols, MARGIN_LEFT, tbl_top, tbl_w, tbl_h)
    table = table_shape.table

    # 列幅（col_widths 指定があれば使用、なければ均等分配）
    col_widths = tbl_data.get("col_widths")
    if col_widths and len(col_widths) == n_cols:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    else:
        col_w = int(tbl_w / n_cols)
        for i in range(n_cols):
            table.columns[i].width = col_w

    # ヘッダー
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["table_hd"]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = C["white"]
            p.font.name = FONT_NAME

    # データ行
    for i, row in enumerate(rows[:14]):
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["table_alt"]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = C["text_dark"]
                p.font.name = FONT_NAME

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


def build_two_column(prs, slide_data: dict, page_num: int, total: int):
    """2カラムスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    half_w = Inches(5.5)
    col_left = slide_data.get("left", {})
    col_right = slide_data.get("right", {})

    for col_data, x_offset in [(col_left, MARGIN_LEFT), (col_right, Inches(7.0))]:
        if not col_data:
            continue
        # サブタイトル
        if col_data.get("title"):
            tb_title = _add_text_box(slide, x_offset, MARGIN_TOP, half_w, Inches(0.5))
            p = tb_title.text_frame.paragraphs[0]
            _set_paragraph(p, col_data["title"], font_size=14, bold=True, color=C["primary"])

        # 箇条書き
        bullets = col_data.get("bullets", [])
        if bullets:
            tb = _add_text_box(slide, x_offset, MARGIN_TOP + Inches(0.6), half_w, Inches(4.5))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, item in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _set_paragraph(p, f"•  {item}", font_size=11, color=C["text_dark"])
                p.space_after = Pt(4)

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


# ── 目次（TOC）スライド ──

def _build_toc(prs, slide_data: dict, page_num: int, total: int):
    """目次スライド（リンクは後から追加）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", "目次"))

    items = slide_data.get("items", [])
    if not items:
        _add_footer(slide, page_num, total, "目次")
        return

    tb = _add_text_box(slide, MARGIN_LEFT + Inches(0.5), MARGIN_TOP + Inches(0.2),
                       CONTENT_W - Inches(1.0), CONTENT_H)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        label = item if isinstance(item, str) else item.get("label", "")
        run.text = label
        run.font.size = Pt(16)
        run.font.name = FONT_NAME
        run.font.color.rgb = C["secondary"]
        p.space_after = Pt(14)

    _add_footer(slide, page_num, total, "目次")


def _link_toc(prs, toc_slide_idx: int, toc_items: list, slide_title_to_prs_idx: dict):
    """目次の各項目にセクションスライドへのハイパーリンクを追加

    slide_title_to_prs_idx: {title: prs.slides のインデックス} — generate() で構築
    """
    toc_slide = prs.slides[toc_slide_idx]

    # TOC のテキストフレームを探す（タイトル帯・フッターを除いて最大のもの）
    text_frames = [s.text_frame for s in toc_slide.shapes if s.has_text_frame]
    if len(text_frames) < 2:
        return
    # 段落数が最も多いテキストフレームが目次本体
    tf = max(text_frames, key=lambda t: len(t.paragraphs))

    # 空段落をスキップ（paragraphs[0] が空の場合がある）
    real_paras = [p for p in tf.paragraphs if p.runs]

    for para, item in zip(real_paras, toc_items):
        target_title = item if isinstance(item, str) else item.get("target", "")
        if target_title not in slide_title_to_prs_idx:
            continue
        target_idx = slide_title_to_prs_idx[target_title]
        if target_idx >= len(prs.slides):
            continue
        target_slide = prs.slides[target_idx]

        for run in para.runs:
            rId = toc_slide.part.relate_to(
                target_slide.part,
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
            )
            rPr = run._r.get_or_add_rPr()
            hlinkClick = rPr.makeelement(
                qn("a:hlinkClick"),
                {qn("r:id"): rId, "action": "ppaction://hlinksldjump"},
            )
            rPr.append(hlinkClick)


# ── ダイアグラム ヘルパー ──

def _side_point(bx, by, bw, bh, side, gap, frac=0.5):
    """ボックスの指定辺の接続点を返す。frac=0.0〜1.0 で辺上の位置を指定（0.5=中央）"""
    if side == "right":
        return bx + bw + gap, by + bh * frac
    elif side == "left":
        return bx - gap,       by + bh * frac
    elif side == "bottom":
        return bx + bw * frac, by + bh + gap
    elif side == "top":
        return bx + bw * frac, by - gap
    return cx, cy


# ── ダイアグラムスライド ──

# ボックスのスタイル定義（背景色, テキスト色）
_BOX_STYLES = {
    "primary":   (C["primary"],   C["white"]),
    "secondary": (C["secondary"], C["white"]),
    "accent":    (C["accent"],    C["white"]),
    "light":     (C["bg_light"],  C["text_dark"]),
    "white":     (C["white"],     C["text_dark"]),
}


def build_diagram(prs, slide_data: dict, page_num: int, total: int):
    """ダイアグラム（ボックス + 矢印 + グループ）スライド

    JSON elements 形式:
    {
      "groups": [{"label": "...", "x": 1.0, "y": 1.5, "w": 5, "h": 3, "style": "light"}],
      "boxes":  [{"id": "a", "label": "...", "x": 2, "y": 2, "w": 2.5, "h": 1, "style": "primary"}],
      "arrows": [{"from": "a", "to": "b", "label": "..."}]
    }
    座標はスライド左上からの絶対値（インチ）。タイトル帯は 0-1.0 インチ。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    elements = slide_data.get("elements", {})
    groups = elements.get("groups", [])
    boxes = elements.get("boxes", [])
    arrows = elements.get("arrows", [])

    # --- グループ（背景領域）---
    for grp in groups:
        x, y, w, h = Inches(grp["x"]), Inches(grp["y"]), Inches(grp["w"]), Inches(grp["h"])
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        bg, _ = _BOX_STYLES.get(grp.get("style", "light"), _BOX_STYLES["light"])
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = C["border"]
        shape.line.width = Pt(1)
        # ラベル（左上に小さく）
        if grp.get("label"):
            tf = shape.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = grp["label"]
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = C["text_sub"]
            tf.paragraphs[0].font.name = FONT_NAME
            tf.paragraphs[0].font.bold = True
            tf.vertical_anchor = MSO_ANCHOR.TOP

    # --- ボックス ---
    box_map = {}  # id → {x, y, w, h}
    for box in boxes:
        bx, by, bw, bh = box["x"], box["y"], box["w"], box["h"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bx), Inches(by), Inches(bw), Inches(bh),
        )
        bg_color, text_color = _BOX_STYLES.get(box.get("style", "primary"), _BOX_STYLES["primary"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = bg_color
        shape.line.width = Pt(0)

        # 影（軽いドロップシャドウ）
        sp = shape._element
        spPr = sp.find(qn("a:spPr"))
        if spPr is None:
            spPr = sp.makeelement(qn("a:spPr"), {})
            sp.append(spPr)
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
            "blurRad": "40000", "dist": "20000", "dir": "5400000",
            "rotWithShape": "0",
        })
        srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
        alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "30000"})
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines = box["label"].split("\n")
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(10)
            p.font.bold = (li == 0)
            p.font.color.rgb = text_color
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER

        box_map[box["id"]] = {"x": bx, "y": by, "w": bw, "h": bh}

    # --- 矢印（シンプル直線）---
    # ※ 矢印が他のボックスと被らないよう、JSON 側で隣接配置を保証すること
    GAP = 0.15  # ボックスの辺からのオフセット（インチ）

    for arr in arrows:
        src = box_map.get(arr.get("from"))
        dst = box_map.get(arr.get("to"))
        if not src or not dst:
            continue

        sx, sy, sw, sh = src["x"], src["y"], src["w"], src["h"]
        dx, dy, dw, dh = dst["x"], dst["y"], dst["w"], dst["h"]
        scx, scy = sx + sw / 2, sy + sh / 2
        dcx, dcy = dx + dw / 2, dy + dh / 2

        # JSON で接続辺を指定できる（省略時は自動判定）
        side_from = arr.get("side_from")  # "top" / "bottom" / "left" / "right"
        side_to   = arr.get("side_to")
        sf_frac   = arr.get("side_from_frac", 0.5)
        st_frac   = arr.get("side_to_frac",   0.5)

        if side_from and side_to:
            x1, y1 = _side_point(sx, sy, sw, sh, side_from, GAP, sf_frac)
            x2, y2 = _side_point(dx, dy, dw, dh, side_to,   GAP, st_frac)
        else:
            # 自動判定（近い辺の中央）
            ddx, ddy = dcx - scx, dcy - scy
            if abs(ddx) > abs(ddy):
                if ddx > 0:
                    x1, y1 = sx + sw + GAP, scy
                    x2, y2 = dx - GAP, dcy
                else:
                    x1, y1 = sx - GAP, scy
                    x2, y2 = dx + dw + GAP, dcy
            else:
                if ddy > 0:
                    x1, y1 = scx, sy + sh + GAP
                    x2, y2 = dcx, dy - GAP
                else:
                    x1, y1 = scx, sy - GAP
                    x2, y2 = dcx, dy + dh + GAP

        connector = slide.shapes.add_connector(
            1,  # STRAIGHT
            Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        if arr.get("arrow_style") == "master_detail":
            connector.line.width = Pt(2.5)
            connector.line.color.rgb = C["primary"]
        else:
            connector.line.width = Pt(1.5)
            connector.line.color.rgb = C["text_sub"]

        # 矢じり（tail = 終端, head = 始端）
        ln = connector.line._ln
        tail = ln.makeelement(qn("a:tailEnd"), {
            "type": "triangle", "w": "med", "len": "med",
        })
        ln.append(tail)
        # 双方向矢印（"bidirectional": true）
        if arr.get("bidirectional"):
            head = ln.makeelement(qn("a:headEnd"), {
                "type": "triangle", "w": "med", "len": "med",
            })
            ln.append(head)

        # ラベル（矢印の中間点）
        if arr.get("label"):
            lx = (x1 + x2) / 2
            ly = (y1 + y2) / 2
            lbl_w = max(len(arr["label"]) * 0.16, 0.8)
            lbl_h = 0.26

            # 水平矢印ならラベルを上にオフセット、垂直矢印なら右にオフセット
            if abs(x2 - x1) > abs(y2 - y1):
                ly -= 0.2
            else:
                lx += 0.2

            lbl_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(lx - lbl_w / 2), Inches(ly - lbl_h / 2),
                Inches(lbl_w), Inches(lbl_h),
            )
            lbl_shape.fill.solid()
            lbl_shape.fill.fore_color.rgb = C["white"]
            lbl_shape.line.color.rgb = C["border"]
            lbl_shape.line.width = Pt(0.5)
            tf = lbl_shape.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            _set_paragraph(p, arr["label"], font_size=7, color=C["text_sub"], align=PP_ALIGN.CENTER)

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


def build_mermaid_diagram(prs, slide_data: dict, page_num: int, total: int):
    """Mermaid テキストを mmdc で PNG レンダリングし、スライドに埋め込む

    JSON 形式:
    {
      "layout": "mermaid",
      "title": "スライドタイトル",
      "mermaid": "flowchart LR\\n  A[開始] --> B[処理] --> C[終了]",
      "caption": "図の補足説明（省略可）"
    }
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    mermaid_text = slide_data.get("mermaid", "")
    if not mermaid_text:
        _add_footer(slide, page_num, total, slide_data.get("title", ""))
        return

    # Mermaid テーマ設定（パレットに合わせた配色）
    mermaid_config = {
        "theme": "base",
        "themeVariables": {
            "primaryColor": "#4A90D9",
            "primaryTextColor": "#FFFFFF",
            "primaryBorderColor": "#1E3A5F",
            "secondaryColor": "#1E3A5F",
            "secondaryTextColor": "#FFFFFF",
            "secondaryBorderColor": "#1E3A5F",
            "tertiaryColor": "#2C6FAC",
            "tertiaryTextColor": "#FFFFFF",
            "tertiaryBorderColor": "#1E3A5F",
            "lineColor": "#6B6B6B",
            "textColor": "#2D2D2D",
            "edgeLabelBackground": "#FFFFFF",
            "clusterBkg": "#EBF4FF",
            "clusterBorder": "#4A90D9",
            "fontFamily": "sans-serif",
            "fontSize": "16px",
        }
    }

    with tempfile.TemporaryDirectory() as tmpdir:
        mmd_path = os.path.join(tmpdir, "diagram.mmd")
        png_path = os.path.join(tmpdir, "diagram.png")
        cfg_path = os.path.join(tmpdir, "config.json")

        Path(mmd_path).write_text(mermaid_text, encoding="utf-8")
        Path(cfg_path).write_text(json.dumps(mermaid_config, ensure_ascii=False), encoding="utf-8")

        # Windows では mmdc.cmd を使用（shutil.which で自動解決）
        mmdc_cmd = shutil.which("mmdc") or shutil.which("mmdc.cmd") or "mmdc"
        result = subprocess.run(
            f'"{mmdc_cmd}" -i "{mmd_path}" -o "{png_path}" -c "{cfg_path}" -w 1600 --backgroundColor white',
            capture_output=True, text=True, shell=True
        )

        if result.returncode != 0:
            print(f"Mermaid レンダリングエラー: {result.stderr}", file=sys.stderr)
            # エラーテキストをスライドに表示
            tb = _add_text_box(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            _set_paragraph(p, f"図のレンダリングに失敗しました: {result.stderr[:100]}", font_size=11, color=C["accent"])
            _add_footer(slide, page_num, total, slide_data.get("title", ""))
            return

        # caption がある場合は図のエリアを少し縮める
        caption = slide_data.get("caption", "")
        img_h = CONTENT_H - Inches(0.5) if caption else CONTENT_H - Inches(0.2)

        # PNG をスライドに配置（幅固定・縦は比率維持）
        pic = slide.shapes.add_picture(png_path, MARGIN_LEFT, MARGIN_TOP, width=CONTENT_W)

        # 高さがエリアをはみ出す場合は高さ基準に縮小
        if pic.height > img_h:
            ratio = img_h / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = img_h
            # 中央寄せ
            pic.left = int((SLIDE_W - pic.width) / 2)

        # caption
        if caption:
            cap_top = MARGIN_TOP + img_h + Inches(0.1)
            tb = _add_text_box(slide, MARGIN_LEFT, cap_top, CONTENT_W, Inches(0.4))
            p = tb.text_frame.paragraphs[0]
            _set_paragraph(p, caption, font_size=10, color=C["text_sub"], align=PP_ALIGN.CENTER)

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


def build_swimlane(prs, slide_data: dict, page_num: int, total: int):
    """スイムレーン図スライド

    各アクター・システムを横レーンで区切り、処理ステップを列順に並べる。
    矢印がレーン間を横断して「誰から誰へ」のフローを明示する。

    JSON elements 形式:
    {
      "lanes": [{"id": "cust", "label": "顧客", "style": "secondary"}, ...],
      "steps": [{"id": "s1", "lane": "cust", "col": 1, "num": 1,
                 "label": "問い合わせ\\nLPフォーム", "style": "white"}, ...],
      "arrows": [{"from": "s1", "to": "s2", "label": "任意"}, ...]
    }
    col は 1 始まりの整数。num はステップ番号（省略可）。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    elements  = slide_data.get("elements", {})
    lanes     = elements.get("lanes", [])
    steps     = elements.get("steps", [])
    arrows    = elements.get("arrows", [])

    if not lanes or not steps:
        _add_footer(slide, page_num, total, slide_data.get("title", ""))
        return

    # ── レイアウト定数 ──
    LABEL_W  = Inches(1.3)
    X0       = Inches(0.5)                     # 左端（レーンラベル開始）
    X_END    = SLIDE_W - Inches(0.3)
    Y0       = MARGIN_TOP - Inches(0.1)        # レーン上端
    Y_END    = SLIDE_H   - Inches(0.5)

    proc_x0  = X0 + LABEL_W                   # ステップエリア左端（EMU）
    proc_w   = X_END - proc_x0                # ステップエリア幅（EMU）

    n_lanes  = len(lanes)
    n_cols   = max((s.get("col", 1) for s in steps), default=1)

    lane_h   = int((Y_END - Y0) / n_lanes)    # 1レーンの高さ（EMU）
    col_w    = int(proc_w / n_cols)            # 1列の幅（EMU）

    BOX_W    = int(col_w  * 0.78)
    BOX_H    = int(lane_h * 0.62)

    # 列数に応じてステップラベルのフォントサイズを縮小
    STEP_FONT = 8 if n_cols <= 5 else 7 if n_cols <= 7 else 6

    # レーンカラーパレット: style → (ラベル背景, ラベル文字, 行背景)
    _LANE_COLORS = {
        "primary":   (C["primary"],   C["white"],      RGBColor(0xEB, 0xEF, 0xF8)),
        "secondary": (C["secondary"], C["white"],      RGBColor(0xEB, 0xF4, 0xFF)),
        "accent":    (C["accent"],    C["white"],      RGBColor(0xFD, 0xF0, 0xE6)),
        "light":     (C["bg_light"],  C["text_dark"],  RGBColor(0xF5, 0xF7, 0xFA)),
    }

    # ── レーン描画 ──
    lane_meta = {}  # id → {"y_top": EMU, "cy": EMU, "label_color": RGB}

    for i, lane in enumerate(lanes):
        ly_top = Y0 + i * lane_h
        style  = lane.get("style", "light")
        lbl_bg, lbl_fg, row_bg = _LANE_COLORS.get(style, _LANE_COLORS["light"])

        lane_meta[lane["id"]] = {
            "y_top":       ly_top,
            "cy":          ly_top + lane_h // 2,
            "label_color": lbl_bg,
        }

        # 行背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, X0, ly_top, X_END - X0, lane_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = row_bg
        bg.line.fill.background()

        # 区切り線
        sep = slide.shapes.add_connector(1, X0, ly_top, X_END, ly_top)
        sep.line.color.rgb = C["border"]
        sep.line.width = Pt(0.75)

        # ラベル背景
        lbl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, X0, ly_top, LABEL_W, lane_h)
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = lbl_bg
        lbl.line.fill.background()

        # ラベルテキスト
        tf = lbl.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = lane.get("label", "")
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = lbl_fg
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # 最下段区切り線
    slide.shapes.add_connector(1, X0, Y_END, X_END, Y_END).line.color.rgb = C["border"]

    # ── ステップボックス描画 ──
    step_map = {}  # id → (cx, cy, bw, bh)  ※ EMU

    for step in steps:
        lid  = step.get("lane")
        col  = step.get("col", 1)
        meta = lane_meta.get(lid)
        if not meta:
            continue

        cx = proc_x0 + (col - 1) * col_w + col_w // 2
        cy = meta["cy"]
        bx = cx - BOX_W // 2
        by = cy - BOX_H // 2

        style    = step.get("style", "white")
        bg_color, text_color = _BOX_STYLES.get(style, _BOX_STYLES["white"])

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, BOX_W, BOX_H)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = meta["label_color"]
        shape.line.width = Pt(1.0)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 列数に応じた1行あたりの表示幅上限（全角換算）
        _max_w = 14 if n_cols >= 8 else 18 if n_cols >= 6 else 22
        raw_lines = step.get("label", "").split("\n")
        lines = [_swimlane_trunc(ln, _max_w) for ln in raw_lines[:4]]
        if len(raw_lines) > 4:
            lines.append("…")
        for li, ln_text in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.text = ln_text
            p.font.size = Pt(STEP_FONT)
            p.font.bold = (li == 0)
            p.font.color.rgb = text_color
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER

        # テキストが枠に収まるよう normAutofit を設定（フォント自動縮小）
        txBody  = tf._txBody
        bodyPr  = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            for _tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
                _el = bodyPr.find(qn(_tag))
                if _el is not None:
                    bodyPr.remove(_el)
            bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))

        # ステップ番号バッジ（左上）
        num = step.get("num")
        if num is not None:
            NW = Inches(0.22)
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, bx, by, NW, NW)
            badge.fill.solid()
            badge.fill.fore_color.rgb = meta["label_color"]
            badge.line.fill.background()
            np_ = badge.text_frame.paragraphs[0]
            np_.text = str(num)
            np_.font.size = Pt(7)
            np_.font.bold = True
            np_.font.color.rgb = C["white"]
            np_.font.name = FONT_NAME_EN
            np_.alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        step_map[step["id"]] = (cx, cy, BOX_W, BOX_H)

    # ── 矢印描画 ──
    GAP = Inches(0.08)

    for arr in arrows:
        src = step_map.get(arr.get("from"))
        dst = step_map.get(arr.get("to"))
        if not src or not dst:
            continue

        sx, sy, sw, sh = src
        dx, dy, dw, dh = dst

        # 水平成分が大きい（列方向）: 右辺→左辺
        # 垂直成分が大きい（同列別レーン）: 下辺/上辺
        ddx, ddy = dx - sx, dy - sy

        if abs(ddx) >= abs(ddy) or ddx > 0:
            if ddx >= 0:
                x1, y1 = sx + sw // 2 + GAP, sy
                x2, y2 = dx - dw // 2 - GAP, dy
            else:
                x1, y1 = sx - sw // 2 - GAP, sy
                x2, y2 = dx + dw // 2 + GAP, dy
        else:
            if ddy >= 0:
                x1, y1 = sx, sy + sh // 2 + GAP
                x2, y2 = dx, dy - dh // 2 - GAP
            else:
                x1, y1 = sx, sy - sh // 2 - GAP
                x2, y2 = dx, dy + dh // 2 + GAP

        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.width = Pt(1.5)
        conn.line.color.rgb = C["secondary"]

        ln = conn.line._ln
        tail = ln.makeelement(qn("a:tailEnd"), {
            "type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)

        if arr.get("label"):
            lx = (x1 + x2) // 2
            ly = (y1 + y2) // 2
            lw = Inches(max(len(arr["label"]) * 0.13, 0.7))
            lh = Inches(0.26)
            lbl = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                lx - lw // 2, ly - lh // 2, lw, lh)
            lbl.fill.solid()
            lbl.fill.fore_color.rgb = C["white"]
            lbl.line.color.rgb = C["border"]
            lbl.line.width = Pt(0.5)
            p = lbl.text_frame.paragraphs[0]
            _set_paragraph(p, arr["label"], font_size=7,
                           color=C["text_sub"], align=PP_ALIGN.CENTER)

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


# ── ER 図スライド（header + field list per box）──


def _draw_crow_foot_marker(slide, x, y, is_horizontal, toward_box_pos,
                            kind, color, lw=Pt(1.4)):
    """ER線の端点にクロウズフット記法のマーカーを描画する。

    Args:
      x, y: 端点座標（インチ, ボックス側の線端）
      is_horizontal: 本線が水平かどうか
      toward_box_pos: ボックスが端点から正方向（+x or +y）にあるなら True
      kind:
        "one"       : 縦バー1本（必須 one）
        "one_md"    : ダイヤモンド（主従親側・master側）
        "one_opt"   : 小円 + 縦バー（参照親側・optional one）
        "many"      : crow's foot（必須 many）
        "many_opt"  : 小円 + crow's foot（参照子側・optional many）
      color: 線色
    """
    BAR    = 0.075   # 1本線バーの半長
    FORK_D = 0.13    # クロウズフット apex からの奥行き
    SPREAD = 0.085   # フォーク開き幅
    sign = 1 if toward_box_pos else -1

    def _draw_bar(bar_x, bar_y):
        """(bar_x, bar_y) を中心とする垂直（水平線の場合）バー"""
        if is_horizontal:
            line = slide.shapes.add_connector(
                1, Inches(bar_x), Inches(bar_y - BAR),
                   Inches(bar_x), Inches(bar_y + BAR),
            )
        else:
            line = slide.shapes.add_connector(
                1, Inches(bar_x - BAR), Inches(bar_y),
                   Inches(bar_x + BAR), Inches(bar_y),
            )
        line.line.color.rgb = color
        line.line.width = Pt(2.0)

    def _draw_fan(apex_xy, fan_pts):
        """apex から fan の各点へ 3 本の線を引く"""
        for fx, fy in fan_pts:
            line = slide.shapes.add_connector(
                1, Inches(apex_xy[0]), Inches(apex_xy[1]),
                   Inches(fx), Inches(fy),
            )
            line.line.color.rgb = color
            line.line.width = lw

    def _draw_circle(cx_, cy_):
        """(cx_, cy_) を中心とする白塗り小円（径 CIRC_D）"""
        CIRC_D = 0.10
        ox = cx_ - CIRC_D / 2
        oy = cy_ - CIRC_D / 2
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(ox), Inches(oy), Inches(CIRC_D), Inches(CIRC_D),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = C["white"]
        circ.line.color.rgb = color
        circ.line.width = Pt(1.2)

    def _draw_diamond(cx_, cy_):
        """(cx_, cy_) を中心とする小ダイヤモンド（主従親マーカー）"""
        DIA_W, DIA_H = 0.14, 0.10
        ox = cx_ - DIA_W / 2
        oy = cy_ - DIA_H / 2
        dia = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Inches(ox), Inches(oy),
            Inches(DIA_W), Inches(DIA_H),
        )
        dia.fill.solid()
        dia.fill.fore_color.rgb = color
        dia.line.color.rgb = color
        dia.line.width = Pt(0.75)

    # ── 各 kind 実装 ──
    if kind == "one":
        off = -sign * 0.02
        if is_horizontal:
            _draw_bar(x + off, y)
        else:
            _draw_bar(x, y + off)
        return

    if kind == "one_md":
        # 端点の外側（ボックス逆方向）にダイヤモンドを配置
        if is_horizontal:
            _draw_diamond(x - sign * (FORK_D * 0.7), y)
        else:
            _draw_diamond(x, y - sign * (FORK_D * 0.7))
        return

    if kind == "one_opt":
        # 外側（線側）に円、内側（ボックス寄り）にバー
        if is_horizontal:
            circ_x = x - sign * (FORK_D + 0.05)
            bar_x  = x - sign * 0.02
            _draw_circle(circ_x, y)
            _draw_bar(bar_x, y)
        else:
            circ_y = y - sign * (FORK_D + 0.05)
            bar_y  = y - sign * 0.02
            _draw_circle(x, circ_y)
            _draw_bar(x, bar_y)
        return

    if kind == "many":
        if is_horizontal:
            apex = (x - sign * FORK_D, y)
            fan = [(x, y - SPREAD), (x, y), (x, y + SPREAD)]
        else:
            apex = (x, y - sign * FORK_D)
            fan = [(x - SPREAD, y), (x, y), (x + SPREAD, y)]
        _draw_fan(apex, fan)
        return

    if kind == "many_opt":
        # 外側に円、内側に crow's foot
        if is_horizontal:
            circ_x = x - sign * (FORK_D + 0.12)
            _draw_circle(circ_x, y)
            apex = (x - sign * FORK_D, y)
            fan = [(x, y - SPREAD), (x, y), (x, y + SPREAD)]
        else:
            circ_y = y - sign * (FORK_D + 0.12)
            _draw_circle(x, circ_y)
            apex = (x, y - sign * FORK_D)
            fan = [(x - SPREAD, y), (x, y), (x + SPREAD, y)]
        _draw_fan(apex, fan)
        return



_FIELD_TYPE_ICON = {
    "reference": "↗",
    "id":        "🔑",
    "picklist":  "▼",
    "multipicklist": "▼",
    "boolean":   "☑",
    "currency":  "¥",
    "number":    "#",
    "percent":   "%",
    "date":      "D",
    "datetime":  "D",
    "formula":   "fx",
    "autonumber":"#",
}

_ER_HDR_COLORS = {
    # style → (hdr_bg, hdr_fg, body_bg, body_fg, border)
    "primary":   (C["primary"],   C["white"], RGBColor(0xE8, 0xF0, 0xF8), C["text_dark"], C["primary"]),
    "secondary": (C["secondary"], C["white"], RGBColor(0xE8, 0xF4, 0xFF), C["text_dark"], C["secondary"]),
    "accent":    (C["accent"],    C["white"], RGBColor(0xFF, 0xF2, 0xE8), C["text_dark"], C["accent"]),
    "light":     (C["bg_light"],  C["text_dark"], C["white"], C["text_dark"], C["border"]),
    "white":     (C["white"],     C["text_dark"], C["white"], C["text_dark"], C["border"]),
    "ref":       (RGBColor(0xCC, 0xCC, 0xCC), C["text_dark"],
                  RGBColor(0xF5, 0xF5, 0xF5), C["text_sub"],
                  RGBColor(0xAA, 0xAA, 0xAA)),
}


# OWDバッジの配色
_OWD_BADGE_COLORS = {
    "Private":            (RGBColor(0xFF, 0xDD, 0xDD), RGBColor(0x8B, 0x00, 0x00)),
    "ReadWrite":          (RGBColor(0xDD, 0xFF, 0xDD), RGBColor(0x00, 0x66, 0x00)),
    "Public Read Write":  (RGBColor(0xDD, 0xFF, 0xDD), RGBColor(0x00, 0x66, 0x00)),
    "Read Only":          (RGBColor(0xFF, 0xFF, 0xDD), RGBColor(0x80, 0x60, 0x00)),
    "Public Read Only":   (RGBColor(0xFF, 0xFF, 0xDD), RGBColor(0x80, 0x60, 0x00)),
    "ControlledByParent": (RGBColor(0xDD, 0xEE, 0xFF), RGBColor(0x00, 0x40, 0x80)),
}


def _owd_badge_color(owd: str):
    """OWD文字列から (bg, fg) 色を決定する"""
    if not owd:
        return None
    for key, colors in _OWD_BADGE_COLORS.items():
        if key.lower() in owd.lower():
            return colors
    # その他はグレー
    return (RGBColor(0xE0, 0xE0, 0xE0), RGBColor(0x40, 0x40, 0x40))


def _owd_badge_short(owd: str) -> str:
    """OWD表示用の短縮ラベル"""
    if not owd:
        return ""
    low = owd.lower()
    if "controlledbyparent" in low or "親" in owd:
        return "CbP"
    if "private" in low or "非公開" in owd:
        return "Private"
    if "read" in low and "write" in low:
        return "R/W"
    if "read" in low:
        return "R/O"
    return owd[:10]


def _set_shape_dashed(shape, preset="dash"):
    """shape の罫線を破線にする"""
    try:
        ln = shape.line._get_or_add_ln()
    except Exception:
        ln = shape.line._ln
    if ln is None:
        return
    for old in ln.findall(qn('a:prstDash')):
        ln.remove(old)
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    prstDash = parse_xml(f'<a:prstDash xmlns:a="{ns}" val="{preset}"/>')
    ln.append(prstDash)


def _set_connector_dashed(connector, preset="lgDash"):
    """コネクタ線を破線にする"""
    ln = connector.line._ln
    if ln is None:
        return
    for old in ln.findall(qn('a:prstDash')):
        ln.remove(old)
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    prstDash = parse_xml(f'<a:prstDash xmlns:a="{ns}" val="{preset}"/>')
    ln.append(prstDash)

HDR_H = Inches(0.52)    # ヘッダー部分の固定高さ
FIELD_H = Inches(0.215) # 1フィールド行の高さ


def _add_er_box(slide, box: dict) -> tuple:
    """ER図の1オブジェクトボックスを描画する。
    box keys: id, label, api_name, fields, x, y, w, h, style
              + ref_only(bool), owd(str), record_count(str)
    Returns: (x_emu, y_emu, w_emu, h_emu)
    """
    bx, by, bw, bh = box["x"], box["y"], box["w"], box["h"]
    bx_e, by_e = Inches(bx), Inches(by)
    bw_e, bh_e = Inches(bw), Inches(bh)
    hdr_bg, hdr_fg, body_bg, body_fg, border = _ER_HDR_COLORS.get(
        box.get("style", "primary"), _ER_HDR_COLORS["primary"]
    )

    ref_only = bool(box.get("ref_only", False))

    # 1. 全体ボックス（角丸、ボディ色）
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx_e, by_e, bw_e, bh_e)
    outer.fill.solid()
    outer.fill.fore_color.rgb = body_bg
    outer.line.color.rgb = border
    outer.line.width = Pt(1)
    outer.adjustments[0] = 0.03

    # ref_only のときは外枠を破線に
    if ref_only:
        _set_shape_dashed(outer, "dash")

    # 2. ヘッダー矩形
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx_e, by_e, bw_e, HDR_H)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = hdr_bg
    hdr.line.fill.background()

    # OWDバッジをヘッダー右端に描画
    owd = box.get("owd", "")
    badge_cols = _owd_badge_color(owd)
    badge_w_in = 0.0
    if badge_cols and not ref_only:
        bg_col, fg_col = badge_cols
        label_short = _owd_badge_short(owd)
        badge_w_in = min(max(len(label_short) * 0.07 + 0.18, 0.45), 0.85)
        badge_h_in = 0.22
        badge_x = bx + bw - badge_w_in - 0.08
        badge_y = by + (HDR_H / Inches(1)) / 2 - badge_h_in / 2
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(badge_x), Inches(badge_y),
            Inches(badge_w_in), Inches(badge_h_in),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = bg_col
        badge.line.color.rgb = bg_col
        badge.line.width = Pt(0.25)
        btf = badge.text_frame
        btf.margin_left = Inches(0.02)
        btf.margin_right = Inches(0.02)
        btf.margin_top = Inches(0.0)
        btf.margin_bottom = Inches(0.0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        _set_paragraph(bp, label_short, font_size=6.5, bold=True,
                       color=fg_col, align=PP_ALIGN.CENTER)

    # 3. ヘッダーテキスト（badge 分幅を縮める）
    hdr_text_w_margin = 0.14 + (badge_w_in + 0.06 if badge_w_in else 0)
    hdr_tb = slide.shapes.add_textbox(
        bx_e + Inches(0.07), by_e + Inches(0.04),
        bw_e - Inches(hdr_text_w_margin), HDR_H - Inches(0.04),
    )
    hdr_tf = hdr_tb.text_frame
    hdr_tf.word_wrap = True

    p0 = hdr_tf.paragraphs[0]
    p0.text = box.get("label", box.get("api_name", ""))
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = hdr_fg
    p0.font.name = FONT_NAME
    p0.space_after = Pt(0)

    api_name = box.get("api_name", "")
    if api_name and api_name != box.get("label", ""):
        p1 = hdr_tf.add_paragraph()
        p1.text = api_name
        p1.font.size = Pt(7)
        p1.font.bold = False
        p1.font.color.rgb = hdr_fg
        p1.font.name = FONT_NAME_EN
        p1.space_after = Pt(0)

    # 4. フィールドリスト（ref_only のときはスキップ）
    fields = [] if ref_only else box.get("fields", [])
    if fields:
        fld_tb = slide.shapes.add_textbox(
            bx_e + Inches(0.07),
            by_e + HDR_H + Inches(0.03),
            bw_e - Inches(0.14),
            bh_e - HDR_H - Inches(0.04),
        )
        fld_tf = fld_tb.text_frame
        fld_tf.word_wrap = True

        for i, fld in enumerate(fields):
            p = fld_tf.paragraphs[0] if i == 0 else fld_tf.add_paragraph()
            icon   = _FIELD_TYPE_ICON.get(fld.get("type", ""), "•")
            label  = fld.get("label") or fld.get("api_name", "")
            api    = fld.get("api_name", "")
            is_fk  = fld.get("is_fk", False) or fld.get("type") in ("reference",)
            ref_lbl = fld.get("ref_label", "")

            if is_fk:
                display = label if label != api else api
                if ref_lbl and ref_lbl not in display:
                    text = f"↗ {display}  →{ref_lbl}"
                else:
                    text = f"↗ {display}"
            else:
                text = f"{icon} {label}"
            p.text = text
            p.font.size = Pt(7)
            p.font.bold = is_fk
            p.font.color.rgb = body_fg if not is_fk else C["primary"] if hdr_bg != C["primary"] else C["text_dark"]
            p.font.name = FONT_NAME
            p.space_after = Pt(1)

    # 5. レコード数をボックス下端に小さく表示
    record_count = box.get("record_count", "")
    if record_count and not ref_only:
        rc_text = record_count if ("件" in record_count or "レコード" in record_count) else f"{record_count}件"
        rc_tb = slide.shapes.add_textbox(
            bx_e + Inches(0.07),
            by_e + bh_e - Inches(0.22),
            bw_e - Inches(0.14),
            Inches(0.18),
        )
        rc_tf = rc_tb.text_frame
        rc_tf.margin_top = Inches(0)
        rc_tf.margin_bottom = Inches(0)
        rcp = rc_tf.paragraphs[0]
        _set_paragraph(rcp, rc_text, font_size=6,
                       color=C["text_sub"], align=PP_ALIGN.RIGHT)

    return bx_e, by_e, bw_e, bh_e


def build_er(prs, slide_data: dict, page_num: int, total: int):
    """ER図スライド（オブジェクト名 + フィールドリスト + リレーション矢印）

    JSON elements 形式:
    {
      "boxes": [{
        "id": "Account", "label": "取引先", "api_name": "Account",
        "fields": [{"api_name": "Name", "type": "string", "label": "取引先名"}, ...],
        "x": 1.0, "y": 1.5, "w": 2.5, "h": 1.6, "style": "secondary"
      }],
      "arrows": [{"from": "A", "to": "B", "label": "1対多", "arrow_style": "master_detail"}]
    }
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", ""))

    elements = slide_data.get("elements", {})
    boxes    = elements.get("boxes", [])
    arrows   = elements.get("arrows", [])

    # --- ボックス描画 ---
    box_map: dict[str, dict] = {}
    for box in boxes:
        _add_er_box(slide, box)
        box_map[box["id"]] = {"x": box["x"], "y": box["y"],
                               "w": box["w"], "h": box["h"]}

    # --- 矢印（diagram と同じロジック）---
    GAP = 0.12
    for arr in arrows:
        src = box_map.get(arr.get("from"))
        dst = box_map.get(arr.get("to"))
        if not src or not dst:
            continue

        sx, sy, sw, sh = src["x"], src["y"], src["w"], src["h"]
        dx, dy, dw, dh = dst["x"], dst["y"], dst["w"], dst["h"]
        scx, scy = sx + sw / 2, sy + sh / 2
        dcx, dcy = dx + dw / 2, dy + dh / 2

        side_from = arr.get("side_from")
        side_to   = arr.get("side_to")
        sf_frac   = arr.get("side_from_frac", 0.5)
        st_frac   = arr.get("side_to_frac",   0.5)
        if side_from and side_to:
            x1, y1 = _side_point(sx, sy, sw, sh, side_from, GAP, sf_frac)
            x2, y2 = _side_point(dx, dy, dw, dh, side_to,   GAP, st_frac)
        else:
            ddx, ddy = dcx - scx, dcy - scy
            if abs(ddx) >= abs(ddy):
                if ddx > 0:
                    x1, y1 = sx + sw + GAP, scy
                    x2, y2 = dx - GAP, dcy
                else:
                    x1, y1 = sx - GAP, scy
                    x2, y2 = dx + dw + GAP, dcy
            else:
                if ddy > 0:
                    x1, y1 = scx, sy + sh + GAP
                    x2, y2 = dcx, dy - GAP
                else:
                    x1, y1 = scx, sy - GAP
                    x2, y2 = dcx, dy + dh + GAP

        connector = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        is_md = arr.get("arrow_style") == "master_detail"
        line_color = RGBColor(0xCC, 0x22, 0x00) if is_md else RGBColor(0x2C, 0x6F, 0xAC)
        connector.line.width = Pt(2.5 if is_md else 1.5)
        connector.line.color.rgb = line_color
        # Lookup は破線に
        if not is_md:
            _set_connector_dashed(connector, "lgDash")

        # マーカー: MD = ダイヤモンド + crow's foot / Lookup = 円+バー + 円+crow's foot
        is_horiz = abs(x2 - x1) >= abs(y2 - y1)
        parent_toward_pos = (x1 < scx) if is_horiz else (y1 < scy)
        child_toward_pos  = (x2 < dcx) if is_horiz else (y2 < dcy)
        if is_md:
            _draw_crow_foot_marker(slide, x1, y1, is_horiz, parent_toward_pos,
                                    "one_md", line_color)
            _draw_crow_foot_marker(slide, x2, y2, is_horiz, child_toward_pos,
                                    "many", line_color)
        else:
            _draw_crow_foot_marker(slide, x1, y1, is_horiz, parent_toward_pos,
                                    "one_opt", line_color)
            _draw_crow_foot_marker(slide, x2, y2, is_horiz, child_toward_pos,
                                    "many_opt", line_color)

        if arr.get("bidirectional"):
            _draw_crow_foot_marker(slide, x1, y1, is_horiz, parent_toward_pos,
                                    "many", line_color)

        # 矢印ラベルは描画しない（マーカーで種別が明示されるため）

    _add_footer(slide, page_num, total, slide_data.get("title", ""))


def build_er_legend(prs, slide_data: dict, page_num: int, total: int):
    """ER図 凡例スライド: リレーション種別・記号・ボックスカラー・OWDバッジの説明"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, slide_data.get("title", "ER図 凡例"))

    MD_COLOR = RGBColor(0xCC, 0x22, 0x00)
    LU_COLOR = RGBColor(0x2C, 0x6F, 0xAC)

    # セクション1: リレーション種別 -------------------------
    y_cursor = 1.1

    def _section_title(text, y):
        tb = _add_text_box(slide, Inches(0.6), Inches(y), Inches(12), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        _set_paragraph(p, text, font_size=13, bold=True, color=C["primary"])

    _section_title("【リレーション種別】", y_cursor)
    y_cursor += 0.40

    # MD サンプル線
    md_y = y_cursor + 0.13
    md_x1, md_x2 = 1.2, 3.6
    md_conn = slide.shapes.add_connector(
        1, Inches(md_x1), Inches(md_y), Inches(md_x2), Inches(md_y),
    )
    md_conn.line.color.rgb = MD_COLOR
    md_conn.line.width = Pt(2.5)
    _draw_crow_foot_marker(slide, md_x1, md_y, True, False, "one_md", MD_COLOR)
    _draw_crow_foot_marker(slide, md_x2, md_y, True, True,  "many",   MD_COLOR)
    tb = _add_text_box(slide, Inches(4.0), Inches(y_cursor), Inches(9), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    _set_paragraph(p,
        "主従（Master-Detail）: 実線・赤  ◆が親（Master）側・|<が子（必須多）",
        font_size=11, color=C["text_dark"])
    y_cursor += 0.40

    # Lookup サンプル線
    lu_y = y_cursor + 0.13
    lu_conn = slide.shapes.add_connector(
        1, Inches(md_x1), Inches(lu_y), Inches(md_x2), Inches(lu_y),
    )
    lu_conn.line.color.rgb = LU_COLOR
    lu_conn.line.width = Pt(1.5)
    _set_connector_dashed(lu_conn, "lgDash")
    _draw_crow_foot_marker(slide, md_x1, lu_y, True, False, "one_opt",  LU_COLOR)
    _draw_crow_foot_marker(slide, md_x2, lu_y, True, True,  "many_opt", LU_COLOR)
    tb = _add_text_box(slide, Inches(4.0), Inches(y_cursor), Inches(9), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    _set_paragraph(p,
        "参照（Lookup）: 破線・青  ○|が親（任意1）・○|<が子（任意多）",
        font_size=11, color=C["text_dark"])
    y_cursor += 0.48

    # セクション2: 記号説明 -------------------------
    _section_title("【記号説明】", y_cursor)
    y_cursor += 0.35

    symbol_rows = [
        ("◆",    "主従の親側（Master）"),
        ("|",    "必須（1以上）"),
        ("○",    "任意（0以上）"),
        ("|<",   "多（crow's foot・必須多）"),
        ("○|<",  "任意の多（0以上）"),
    ]
    for sym, desc in symbol_rows:
        tb = _add_text_box(slide, Inches(0.9), Inches(y_cursor), Inches(1.2), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        _set_paragraph(p, sym, font_size=13, bold=True, color=C["primary"])
        tb2 = _add_text_box(slide, Inches(2.2), Inches(y_cursor), Inches(8), Inches(0.3))
        p2 = tb2.text_frame.paragraphs[0]
        _set_paragraph(p2, desc, font_size=10, color=C["text_dark"])
        y_cursor += 0.25

    y_cursor += 0.12

    # セクション3: ボックスカラー凡例 -------------------------
    _section_title("【ボックスカラー凡例】", y_cursor)
    y_cursor += 0.35

    color_rows = [
        ("primary",   "TX系（トランザクション）オブジェクト"),
        ("accent",    "マスタ系オブジェクト"),
        ("secondary", "標準オブジェクト（外部参照）"),
        ("ref",       "ref_only（破線枠）: TXから参照される外部オブジェクト"),
    ]
    for style, desc in color_rows:
        hdr_bg, _, _, _, _ = _ER_HDR_COLORS.get(style, _ER_HDR_COLORS["light"])
        sw = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(y_cursor), Inches(0.35), Inches(0.22),
        )
        sw.fill.solid()
        sw.fill.fore_color.rgb = hdr_bg
        sw.line.color.rgb = C["border"]
        if style == "ref":
            _set_shape_dashed(sw, "dash")
        tb = _add_text_box(slide, Inches(1.4), Inches(y_cursor - 0.03), Inches(10), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        _set_paragraph(p, desc, font_size=10, color=C["text_dark"])
        y_cursor += 0.26

    y_cursor += 0.10

    # セクション4: OWDバッジ -------------------------
    _section_title("【OWDバッジ】", y_cursor)
    y_cursor += 0.35

    owd_samples = [
        ("Private",            "Private",            "非公開（所有者のみ）"),
        ("ReadWrite",          "R/W",                "参照/更新可能（Public Read/Write）"),
        ("Read Only",          "R/O",                "参照のみ可能（Public Read Only）"),
        ("ControlledByParent", "CbP",                "親オブジェクトに準拠"),
    ]
    bx = 0.9
    for owd, short, desc in owd_samples:
        colors = _owd_badge_color(owd)
        if not colors:
            continue
        bg_col, fg_col = colors
        badge_w = 0.6
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bx), Inches(y_cursor), Inches(badge_w), Inches(0.24),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = bg_col
        badge.line.color.rgb = bg_col
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        btf.margin_top = Inches(0)
        btf.margin_bottom = Inches(0)
        bp = btf.paragraphs[0]
        _set_paragraph(bp, short, font_size=7, bold=True,
                       color=fg_col, align=PP_ALIGN.CENTER)
        tb = _add_text_box(slide, Inches(bx + badge_w + 0.15),
                           Inches(y_cursor - 0.02), Inches(9), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        _set_paragraph(p, f"{owd} — {desc}", font_size=10, color=C["text_dark"])
        y_cursor += 0.26

    _add_footer(slide, page_num, total, slide_data.get("title", "ER図 凡例"))


def build_er_image(prs, slide_data: dict, page_num: int, total: int):
    """ER図を matplotlib で PNG レンダリングしてスライドに埋め込む。
    build_er の画像版。JSON 形式は build_er と同じ。"""
    import tempfile
    try:
        from er_utils import generate_er_image
    except ImportError:
        import sys
        sys.path.insert(0, str(Path(__file__).parent))
        from er_utils import generate_er_image

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    elements = slide_data.get("elements", {})
    boxes    = elements.get("boxes", [])
    arrows   = elements.get("arrows", [])
    title    = slide_data.get("title", "")

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        tmp_png = f.name

    try:
        ok = generate_er_image(boxes, arrows, tmp_png, title=title)
        if ok:
            slide.shapes.add_picture(
                tmp_png, Emu(0), Emu(0), width=SLIDE_W, height=SLIDE_H,
            )
        else:
            tb = _add_text_box(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            _set_paragraph(p, "ER図の生成に失敗しました（matplotlib 未インストール）",
                           font_size=11, color=C["accent"])
    finally:
        try:
            os.unlink(tmp_png)
        except Exception:
            pass

    _add_footer(slide, page_num, total, title)


def build_er_legend_image(prs, slide_data: dict, page_num: int, total: int):
    """ER図凡例を matplotlib で PNG レンダリングしてスライドに埋め込む。"""
    import tempfile
    try:
        from er_utils import generate_er_legend_image
    except ImportError:
        import sys
        sys.path.insert(0, str(Path(__file__).parent))
        from er_utils import generate_er_legend_image

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        tmp_png = f.name

    try:
        ok = generate_er_legend_image(tmp_png)
        if ok:
            slide.shapes.add_picture(
                tmp_png, Emu(0), Emu(0), width=SLIDE_W, height=SLIDE_H,
            )
        else:
            build_er_legend(prs, slide_data, page_num, total)
            return
    finally:
        try:
            os.unlink(tmp_png)
        except Exception:
            pass

    _add_footer(slide, page_num, total, slide_data.get("title", "ER図 凡例"))


# ── ビルダーマップ ──
BUILDERS = {
    "section":    build_section,
    "content":    build_content,
    "bullets":    build_bullets,
    "table":      build_table,
    "two_column": build_two_column,
    "diagram":    build_diagram,
    "er":         build_er_image,        # matplotlib PNG 版（旧 pptx 図形版から置き換え）
    "er_legend":  build_er_legend_image, # matplotlib PNG 版
    "swimlane":   build_swimlane,
    "mermaid":    build_mermaid_diagram,
}


def generate(data: dict, output_path: str):
    """JSON データから PowerPoint を生成"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 表紙
    build_cover(prs, data)

    slides = data.get("slides", [])
    total = len(slides) + 1  # +1 for cover

    toc_slide_idx = None
    toc_items = None

    # スライドタイトル → prs.slides のインデックス（正確なマッピング）
    slide_title_to_prs_idx = {}

    for i, sd in enumerate(slides):
        layout = sd.get("layout", "content")
        prs_idx = len(prs.slides)  # このスライドが追加される位置

        # タイトルを記録
        title = sd.get("title", "")
        if title:
            slide_title_to_prs_idx[title] = prs_idx

        if layout == "toc":
            toc_slide_idx = prs_idx
            toc_items = sd.get("items", [])
            _build_toc(prs, sd, i + 2, total)
        else:
            builder = BUILDERS.get(layout, build_content)
            builder(prs, sd, i + 2, total)

    # 目次のハイパーリンクを追加（全スライド生成後）
    if toc_slide_idx is not None and toc_items:
        _link_toc(prs, toc_slide_idx, toc_items, slide_title_to_prs_idx)

    prs.save(output_path)
    print(f"\n完了: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="JSON → PowerPoint 生成")
    parser.add_argument("--json-file", required=True, help="スライド構成 JSON ファイルのパス")
    parser.add_argument("--output", required=True, help="出力 .pptx ファイルパス")
    args = parser.parse_args()

    json_data = json.loads(Path(args.json_file).read_text(encoding="utf-8"))
    generate(json_data, args.output)


if __name__ == "__main__":
    main()
